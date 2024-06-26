from pptx import Presentation
from pptx.util import Inches, Pt
import os
import ast

def clean_input(final_prompts):
    final_prompts_lists = [ast.literal_eval(prompts) for prompts in final_prompts]
    flat_prompts = [item for sublist in final_prompts_lists for item in sublist]
    individual_prompts = [prompt.strip() for sublist in flat_prompts for prompt in sublist.split('\n') if prompt.strip()]
    print("Panels after cleaning" + str(individual_prompts))

    
    print("Final Panels for SDXL" + str(individual_prompts))
    print("Len of Panels " + str(len(individual_prompts)))
    return individual_prompts[1::]

# def clean_input_ppt(text):
#     print(type(text))
#     print(text)
#     # text = """["\n            1. Tom lives with Aunt Polly after parents' death.\n            2. Tom doesn't like school or work, prefers play and adventure.\n 
#     #        3. Aunt Polly angry, assigns fence painting as punishment, later allows others to join in for food."]"""
#     prompts = text.split('\n')
#     print(" p   " + str(prompts))
#     text_prompts_cleaned = []
#     for i in prompts:
#         j = i.strip()
#         # if len(j) > 5:
#         text_prompts_cleaned.append(j)
#         # else:
#         #     continue
    
 
    

def add_slide(prs, image_path, text):
    slide_layout = prs.slide_layouts[6]  # Use the layout for title and content
    slide = prs.slides.add_slide(slide_layout)

    left_inch = (prs.slide_width - Inches(6)) / 2
    top_inch = (prs.slide_height - Inches(5)) / 2

    # Add image
    slide.shapes.add_picture(image_path, left_inch, top_inch, Inches(6), Inches(4.5))

    # Add text box below image
    txBox = slide.shapes.add_textbox(left_inch / 2, top_inch + Inches(4.5), Inches(6), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(24)  # Increase font size to 24 points
    p.alignment = 1  # Center alignment

def ppt_main(texts):
    text_list = clean_input(texts)
    images_dir = 'out'
    images_path = os.listdir(images_dir)
    existing_pptx = r"pixi-plot\template_ppt.pptx"
    prs = Presentation(existing_pptx)
    # reqd_text=texts[0][1:]
    
    for image_path, text in zip(images_path, text_list):
        image_path = os.path.join(images_dir, image_path)
        add_slide(prs, image_path, text)

    prs.save('pixi-plot/output.pptx')
    # print(reqd_text)


# # existing_pptx = "template_ppt.pptx"
# image_paths = ["generated-1.png", "generated-1.png", "generated-1.png"]  # List of image paths
# # texts = "["\n             Sure, here are the three most important points from the summary:\n\n1. Tom Sawyer is a mischievous and free-spirited boy who doesn't like following rules or going to school.\n2. Tom's Aunt Polly is a strict and authoritative figure who punishes him for not going to school, and Tom is forced to paint the fence as punishment.\n3. Tom is able to manipulate his friends and get them to do things for him, like painting the fence, by offering them food in exchange."]" # List of corresponding texts
# texts = """["\n             Sure, here are the three most important points from the summary:\n\n1. Tom Sawyer is a mischievous and free-spirited boy who doesn't like following rules or going to school.\n2. Tom's Aunt Polly is a strict and authoritative figure who punishes him for not going to school, and Tom is forced to paint the fence as punishment.\n3. Tom is able to manipulate his friends and get them to do things for him, like painting the fence, by offering them food in exchange."]"""
# # # reqd_text=texts[1:]
# ppt_main(texts)
